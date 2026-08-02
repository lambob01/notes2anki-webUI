"""Slide-skipping heuristics.

A false positive here is invisible and expensive: the slide never reaches the
vision model, no card is generated, and nothing in the UI says a slide was
dropped. The figure-bearing cases are the ones that used to fail - a diagram
slide has no extractable text, so a text-only check read it as blank.
"""

import pytest

from app.services.document_reader import MIN_VECTOR_DRAWINGS, DocumentReader

is_title_or_blank = DocumentReader._is_title_or_blank
pdf_page_has_visual = DocumentReader._pdf_page_has_visual
has_visual_content = DocumentReader._has_visual_content


class FakePage:
    def __init__(self, images=(), drawings=()):
        self._images = list(images)
        self._drawings = list(drawings)

    def get_images(self):
        return self._images

    def get_drawings(self):
        return self._drawings


class FakeShape:
    def __init__(self, shape_type=None, **kwargs):
        self.shape_type = shape_type
        for key, value in kwargs.items():
            setattr(self, key, value)


class FakeSlide:
    def __init__(self, *shapes):
        self.shapes = list(shapes)


# --- _is_title_or_blank ----------------------------------------------------


def test_blank_text_is_skipped() -> None:
    assert is_title_or_blank("", 4) is True
    assert is_title_or_blank("   ", 4) is True


def test_early_short_slide_is_skipped() -> None:
    assert is_title_or_blank("Lecture 5: Biology", 0) is True


def test_agenda_slide_is_skipped() -> None:
    assert is_title_or_blank("Agenda\nPart one\nPart two", 3) is True


def test_content_slide_is_kept() -> None:
    text = (
        "Glycolysis converts one glucose molecule into two pyruvate molecules, "
        "yielding a net two ATP and two NADH."
    )
    assert is_title_or_blank(text, 5) is False


def test_figure_slide_with_no_text_is_kept() -> None:
    # The regression: a full-page diagram has no text layer at all.
    assert is_title_or_blank("", 4, has_visual=True) is False


def test_figure_slide_early_in_the_deck_is_kept() -> None:
    assert is_title_or_blank("Fig. 1", 0, has_visual=True) is False


def test_figure_beats_the_agenda_marker() -> None:
    assert is_title_or_blank("Overview", 2, has_visual=True) is False


# --- _pdf_page_has_visual --------------------------------------------------


def test_pdf_page_with_raster_image_has_visual() -> None:
    assert pdf_page_has_visual(FakePage(images=[("img",)])) is True


def test_pdf_page_with_vector_figure_has_visual() -> None:
    page = FakePage(drawings=[{}] * MIN_VECTOR_DRAWINGS)
    assert pdf_page_has_visual(page) is True


def test_pdf_page_with_a_rule_is_not_a_figure() -> None:
    # An underline or a header rule is a drawing too; a couple of ops is not
    # enough to call the page a figure.
    assert pdf_page_has_visual(FakePage(drawings=[{}])) is False


def test_pdf_page_that_cannot_be_classified_is_kept() -> None:
    class Exploding:
        def get_images(self):
            raise RuntimeError("damaged xref")

    assert pdf_page_has_visual(Exploding()) is True


# --- _has_visual_content ---------------------------------------------------


def test_text_only_slide_has_no_visual() -> None:
    pytest.importorskip("pptx")
    assert has_visual_content(FakeSlide(FakeShape())) is False


def test_picture_shape_counts_as_visual() -> None:
    pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    assert has_visual_content(FakeSlide(FakeShape(MSO_SHAPE_TYPE.PICTURE))) is True


def test_chart_in_a_placeholder_counts_as_visual() -> None:
    """A chart inside a content placeholder reports shape_type PLACEHOLDER,
    so the shape_type check alone misses it."""
    pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape = FakeShape(MSO_SHAPE_TYPE.PLACEHOLDER, has_chart=True)
    assert has_visual_content(FakeSlide(shape)) is True


def test_picture_placeholder_counts_as_visual() -> None:
    pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape = FakeShape(MSO_SHAPE_TYPE.PLACEHOLDER, image=object())
    assert has_visual_content(FakeSlide(shape)) is True
