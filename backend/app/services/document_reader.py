from __future__ import annotations

import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Optional

# A PDF page carrying at least this many vector drawing ops is treated as
# holding a figure rather than a rule or an underline. Diagrams exported from
# slide decks arrive as vectors, not embedded images, so counting only raster
# images would read a chart-only page as blank and discard it.
MIN_VECTOR_DRAWINGS = 5


class DocumentSlide:
    def __init__(
        self,
        index: int,
        image_bytes: bytes,
        notes: str = "",
        source_filename: str = "",
    ):
        self.index = index
        self.image_bytes = image_bytes
        self.notes = notes
        self.source_filename = source_filename


class DocumentReader:
    def __init__(self, filepath: str):
        self.filepath = filepath
        self.path = Path(filepath)
        self.ext = self.path.suffix.lower()
        self._slides: Optional[list[DocumentSlide]] = None
        self._text_content: Optional[str] = None

    def is_image(self) -> bool:
        return self.ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"}

    def is_pdf(self) -> bool:
        return self.ext == ".pdf"

    def is_pptx(self) -> bool:
        return self.ext == ".pptx"

    def is_docx(self) -> bool:
        return self.ext == ".docx"

    def is_text(self) -> bool:
        return self.ext in {".txt", ".md"}

    def extract_text(self) -> str:
        if self._text_content is not None:
            return self._text_content

        if self.is_text():
            with open(self.filepath, "r", encoding="utf-8", errors="replace") as f:
                self._text_content = f.read()
            return self._text_content

        if self.is_pdf():
            import fitz
            doc = fitz.open(self.filepath)
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            self._text_content = "\n\n".join(pages)
            return self._text_content

        if self.is_docx():
            from docx import Document
            doc = Document(self.filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            self._text_content = "\n".join(paragraphs)
            return self._text_content

        if self.is_pptx():
            from pptx import Presentation
            prs = Presentation(self.filepath)
            texts = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_text.append(text)
                if slide_text:
                    texts.append(f"Slide {slide_idx + 1}: " + " ".join(slide_text))
            self._text_content = "\n\n".join(texts)
            return self._text_content

        self._text_content = ""
        return self._text_content

    def render_slides(self, dpi: int = 150, skip_title_blank: bool = True) -> list[DocumentSlide]:
        if self._slides is not None:
            return self._slides

        if self.is_image():
            with open(self.filepath, "rb") as f:
                img_bytes = f.read()
            self._slides = [
                DocumentSlide(index=0, image_bytes=img_bytes, source_filename=self.path.name)
            ]
            return self._slides

        if self.is_pdf():
            self._slides = self._render_pdf(dpi, skip_title_blank)
        elif self.is_pptx():
            self._slides = self._render_pptx(dpi, skip_title_blank)
        else:
            self._slides = []

        return self._slides

    def _render_pdf(self, dpi: int, skip_title_blank: bool) -> list[DocumentSlide]:
        import fitz
        slides = []
        doc = fitz.open(self.filepath)

        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if skip_title_blank and self._is_title_or_blank(
                text, i, has_visual=self._pdf_page_has_visual(page)
            ):
                continue

            pix = page.get_pixmap(dpi=dpi)
            img_bytes = pix.tobytes("jpeg")
            slides.append(
                DocumentSlide(
                    index=i,
                    image_bytes=img_bytes,
                    source_filename=self.path.name,
                )
            )

        doc.close()
        return slides

    @staticmethod
    def find_libreoffice() -> Optional[str]:
        """Locate the LibreOffice binary, or None if it isn't installed."""
        import shutil

        found = shutil.which("soffice") or shutil.which("libreoffice")
        if found:
            return found
        for candidate in (
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
        ):
            if os.path.exists(candidate):
                return candidate
        return None

    def _pptx_to_pdf(self, out_dir: Path) -> Optional[Path]:
        """Convert PPTX to PDF with LibreOffice, the high-fidelity path.

        Returns None if LibreOffice is missing or the conversion fails, so the
        caller can fall back to the text renderer.
        """
        import subprocess

        soffice = self.find_libreoffice()
        if not soffice:
            return None

        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(out_dir),
                    str(self.path),
                ],
                check=True,
                capture_output=True,
                # A wedged soffice must not hang the whole generation job.
                timeout=180,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
            return None

        pdf_path = out_dir / f"{self.path.stem}.pdf"
        return pdf_path if pdf_path.exists() else None

    def _pptx_slide_meta(self) -> list[tuple[str, str, bool]]:
        """(visible text, speaker notes, has figure) per slide, in slide order."""
        from pptx import Presentation

        meta = []
        for slide in Presentation(self.filepath).slides:
            text = "\n".join(
                para.text
                for shape in slide.shapes
                if shape.has_text_frame
                for para in shape.text_frame.paragraphs
            )
            notes = ""
            # has_notes_slide first: accessing .notes_slide *creates* one.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
            meta.append((text, notes, self._has_visual_content(slide)))
        return meta

    def _render_pptx(self, dpi: int, skip_title_blank: bool) -> list[DocumentSlide]:
        from PIL import Image

        meta = self._pptx_slide_meta()

        with tempfile.TemporaryDirectory(prefix="notes2anki_pptx_") as tmp:
            tmp_dir = Path(tmp)

            # Preferred path: LibreOffice -> PDF -> PyMuPDF raster. This is what
            # makes vision work; the fallback below can only draw text boxes.
            pdf_path = self._pptx_to_pdf(tmp_dir)
            if pdf_path:
                import fitz

                slides = []
                doc = fitz.open(pdf_path)
                for i, page in enumerate(doc):
                    slide_text, notes, has_visual = (
                        meta[i] if i < len(meta) else ("", "", False)
                    )
                    if skip_title_blank and self._is_title_or_blank(
                        slide_text.strip(), i, has_visual=has_visual
                    ):
                        continue
                    slides.append(
                        DocumentSlide(
                            index=i,
                            image_bytes=page.get_pixmap(dpi=dpi).tobytes("jpeg"),
                            notes=notes,
                            source_filename=self.path.name,
                        )
                    )
                doc.close()
                return slides

            # Fallback: no LibreOffice available. Renders text boxes only, so
            # diagrams and equations are lost, but it beats a blank page.
            from pptx import Presentation

            img_dir = tmp_dir / "slides"
            img_dir.mkdir()
            prs = Presentation(self.filepath)
            slides = []

            for i, slide in enumerate(prs.slides):
                slide_text, notes, has_visual = (
                    meta[i] if i < len(meta) else ("", "", False)
                )
                if skip_title_blank and self._is_title_or_blank(
                    slide_text.strip(), i, has_visual=has_visual
                ):
                    continue

                img_path = img_dir / f"slide_{i:04d}.png"
                self._render_slide_shape_thumbnails(
                    slide, img_path, prs.slide_width, prs.slide_height
                )

                if img_path.exists():
                    img_bytes = img_path.read_bytes()
                else:
                    img = Image.new("RGB", (1280, 720), (255, 255, 255))
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG", quality=85)
                    img_bytes = buf.getvalue()

                slides.append(
                    DocumentSlide(
                        index=i,
                        image_bytes=img_bytes,
                        notes=notes,
                        source_filename=self.path.name,
                    )
                )

        return slides

    @staticmethod
    def _render_slide_shape_thumbnails(
        slide, out_path: Path, slide_width: int, slide_height: int
    ):
        """Draw a slide's text and images with PIL.

        Dimensions come from the Presentation, not the Slide - a Slide has no
        slide_width/slide_height, and reading them off it used to raise
        AttributeError into a bare `except`, so every slide silently rendered
        blank.
        """
        try:
            from io import BytesIO

            from PIL import Image, ImageDraw

            if not slide_width or not slide_height:
                return
            # EMU -> pixels, targeting roughly 1920px wide.
            scale = 1920 / slide_width

            img = Image.new(
                "RGB",
                (int(slide_width * scale), int(slide_height * scale)),
                (255, 255, 255),
            )
            draw = ImageDraw.Draw(img)

            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_blob = shape.image.blob
                        shape_img = Image.open(BytesIO(image_blob))
                        shape_img = shape_img.resize(
                            (int(shape.width * scale), int(shape.height * scale))
                        )
                        img.paste(shape_img, (int(shape.left * scale), int(shape.top * scale)))
                    except Exception:
                        pass

                if shape.has_text_frame:
                    text = " ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    if text:
                        x = int(shape.left * scale)
                        y = int(shape.top * scale)
                        draw.rectangle(
                            [x, y, x + int(shape.width * scale), y + int(shape.height * scale)],
                            outline=(200, 200, 200),
                        )
                        draw.text((x + 5, y + 5), text[:100], fill=(0, 0, 0))

            img.save(out_path, "PNG")
        except Exception:
            # Leave out_path absent; the caller substitutes a blank page. Log
            # it rather than swallowing silently - a bare pass here is what hid
            # the slide_width bug and shipped blank images to the vision model.
            logging.getLogger(__name__).warning(
                "PPTX text fallback failed to render a slide", exc_info=True
            )

    @staticmethod
    def _has_visual_content(slide) -> bool:
        """True if a PPTX slide carries a picture, chart, table, or diagram.

        A slide whose content *is* a figure has little or no extractable text,
        so a text-only heuristic reads it as blank and drops it - discarding
        exactly the material the vision model exists to read.

        Pictures and charts sitting in a content placeholder report shape_type
        PLACEHOLDER rather than PICTURE/CHART, so the has_* flags and `.image`
        are checked too; shape_type alone misses them.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return False

        # getattr per name: an older python-pptx missing one member must not
        # take out the whole check.
        visual_types = {
            getattr(MSO_SHAPE_TYPE, name, None)
            for name in (
                "PICTURE",
                "LINKED_PICTURE",
                "CHART",
                "TABLE",
                "GROUP",
                "MEDIA",
                "EMBEDDED_OLE_OBJECT",
                "LINKED_OLE_OBJECT",
                "DIAGRAM",
                "IGX_GRAPHIC",
                "INK",
                "FREEFORM",
            )
        } - {None}

        for shape in getattr(slide, "shapes", []):
            if getattr(shape, "shape_type", None) in visual_types:
                return True
            if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
                return True
            if hasattr(shape, "image"):
                return True
        return False

    @staticmethod
    def _pdf_page_has_visual(page) -> bool:
        """True if a PDF page carries a raster image or enough vector art to
        be a figure. Errs toward True - a page we cannot classify is kept."""
        try:
            if page.get_images():
                return True
            return len(page.get_drawings()) >= MIN_VECTOR_DRAWINGS
        except Exception:
            return True

    @staticmethod
    def _is_title_or_blank(text: str, slide_index: int, has_visual: bool = False) -> bool:
        # A slide with a figure is content no matter how little text it has:
        # the vision model reads the picture, not the text layer.
        if has_visual:
            return False
        if not text.strip():
            return True
        if slide_index <= 1 and len(text.strip()) < 40:
            return True
        lower = text.lower()
        title_markers = ["agenda", "outline", "overview", "syllabus", "table of contents"]
        if any(marker in lower for marker in title_markers) and len(text.strip()) < 200:
            return True
        return False

    def slide_notes(self, slide_index: int) -> str:
        if not self.is_pptx():
            return ""
        try:
            from pptx import Presentation
            prs = Presentation(self.filepath)
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    return slide.notes_slide.notes_text_frame.text
        except Exception:
            pass
        return ""
